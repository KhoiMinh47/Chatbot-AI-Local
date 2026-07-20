#!/usr/bin/env python3
"""Run the Phase 4 HTTP -> MinIO -> RabbitMQ -> worker -> PostgreSQL acceptance flow."""

from __future__ import annotations

import argparse
import asyncio
import hashlib
import json
import subprocess
import time
from collections.abc import Mapping, Sequence
from datetime import UTC, datetime
from pathlib import Path
from typing import cast
from urllib.parse import quote
from uuid import NAMESPACE_URL, UUID, uuid4, uuid5

import httpx2 as httpx
from app.infrastructure.ingestion_settings import IngestionSettings
from app.main import create_app
from docx import Document
from minio import Minio
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pydantic import SecretStr
from sqlalchemy import URL, create_engine, text
from worker.celery_app import create_celery
from worker.settings import BrokerSettings


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def _write_exclusive(path: Path, content: str) -> None:
    with path.open("x", encoding="utf-8") as handle:
        handle.write(content)


def _font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    for candidate in (
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
    ):
        if Path(candidate).is_file():
            return ImageFont.truetype(candidate, size=size)
    return ImageFont.load_default()


def _save_docx(
    path: Path, title: str, paragraphs: Sequence[str], table: list[list[str]] | None = None
) -> None:
    document = Document()
    document.add_heading(title, level=1)
    for paragraph in paragraphs:
        document.add_paragraph(paragraph)
    if table:
        rendered = document.add_table(rows=len(table), cols=len(table[0]))
        for row_index, row in enumerate(table):
            for column_index, value in enumerate(row):
                rendered.cell(row_index, column_index).text = value
    document.save(str(path))


def _convert_pdf(source: Path, destination_dir: Path, profile_dir: Path) -> Path:
    command = [
        "libreoffice",
        "--headless",
        f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
        "--convert-to",
        "pdf",
        "--outdir",
        str(destination_dir),
        str(source),
    ]
    completed = subprocess.run(command, check=False, capture_output=True, text=True, timeout=120)
    output = destination_dir / f"{source.stem}.pdf"
    if completed.returncode != 0 or not output.is_file():
        raise RuntimeError("LibreOffice failed to create an acceptance PDF")
    return output


def _save_pptx(path: Path, title: str, facts: Sequence[tuple[str, str]]) -> None:
    presentation = Presentation()
    for slide_title, body in facts:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_title
        slide.placeholders[1].text = body
    presentation.save(str(path))


def generate_fixtures(directory: Path) -> tuple[Path, ...]:
    """Generate non-sensitive but structurally real office/PDF acceptance files."""

    directory.mkdir(parents=True)
    profile = directory / ".libreoffice-profile"
    profile.mkdir()

    leave_docx = directory / "source_leave_policy.docx"
    _save_docx(
        leave_docx,
        "Annual Leave Policy",
        (
            "Policy code LEAVE-15 applies to full-time staff.",
            "Eligible staff receive exactly 15 days of annual leave each calendar year.",
        ),
    )
    leave_pdf = _convert_pdf(leave_docx, directory, profile)
    leave_pdf.rename(directory / "01_leave_policy.pdf")
    leave_docx.unlink()

    image = Image.new("RGB", (1800, 1200), "white")
    draw = ImageDraw.Draw(image)
    title_font = _font(64)
    body_font = _font(44)
    draw.text((120, 150), "SCANNED SECURITY NOTICE", fill="black", font=title_font)
    draw.text((120, 330), "Policy code OCR-30.", fill="black", font=body_font)
    draw.text(
        (120, 430),
        "Temporary security badges expire after exactly 30 days.",
        fill="black",
        font=body_font,
    )
    image.save(directory / "02_security_scanned.pdf", "PDF", resolution=180)

    table_docx = directory / "source_benefits_table.docx"
    _save_docx(
        table_docx,
        "Employee Benefits Matrix",
        ("The authoritative benefit codes are listed in the table below.",),
        table=[
            ["Code", "Benefit", "Value"],
            ["MEAL-45", "Meal allowance", "45 USD per month"],
            ["WELL-12", "Wellness days", "12 hours per quarter"],
        ],
    )
    table_pdf = _convert_pdf(table_docx, directory, profile)
    table_pdf.rename(directory / "03_benefits_table.pdf")
    table_docx.unlink()

    (directory / "04_finance_limits.csv").write_text(
        "code,category,limit\nTRAVEL-750,Domestic travel,750 USD\nMEAL-40,Client meal,40 USD\n",
        encoding="utf-8",
    )
    _save_pptx(
        directory / "05_architecture.pptx",
        "Architecture Brief",
        (
            ("Vector Storage", "Architecture code VECTOR-2048 uses vectors with 2048 dimensions."),
            ("Queue", "RabbitMQ carries document ingestion jobs asynchronously."),
        ),
    )
    _save_pptx(
        directory / "06_onboarding.pptx",
        "Onboarding Schedule",
        (
            ("Day One", "Onboarding code START-09 begins at exactly 09:00 local time."),
            ("Equipment", "New staff collect equipment from service desk B2."),
        ),
    )
    _save_docx(
        directory / "07_incident_runbook.docx",
        "Incident Runbook",
        (
            "Runbook code P1-05 requires acknowledgement within exactly 5 minutes.",
            "Escalate unresolved priority-one incidents to the duty manager.",
        ),
    )
    (directory / "08_remote_handbook.txt").write_text(
        "Remote Work Handbook\n\nPolicy code REMOTE-3 permits exactly 3 remote days per week.\n",
        encoding="utf-8",
    )
    (directory / "09_retention_policy.md").write_text(
        "# Retention Policy\n\nPolicy code LOG-90 retains audit logs for exactly **90 days**.\n",
        encoding="utf-8",
    )
    (directory / "10_inventory.csv").write_text(
        "code,item,reorder_level\nLAPTOP-20,Laptop,20\nDOCK-12,Docking station,12\n",
        encoding="utf-8",
    )
    for transient in (profile,):
        for child in sorted(transient.rglob("*"), reverse=True):
            if child.is_file() or child.is_symlink():
                child.unlink()
            elif child.is_dir():
                child.rmdir()
        transient.rmdir()
    return tuple(sorted(path for path in directory.iterdir() if path.is_file()))


def write_gold(path: Path) -> None:
    base_samples: tuple[tuple[str, bool, str, list[str]], ...] = (
        (
            "evaluation",
            True,
            "How many annual leave days does LEAVE-15 grant?",
            ["01_leave_policy.pdf"],
        ),
        (
            "calibration",
            True,
            "When does a temporary badge under OCR-30 expire?",
            ["02_security_scanned.pdf"],
        ),
        (
            "calibration",
            True,
            "What monthly meal allowance is listed for MEAL-45?",
            ["03_benefits_table.pdf"],
        ),
        (
            "calibration",
            True,
            "What is the domestic travel limit for TRAVEL-750?",
            ["04_finance_limits.csv"],
        ),
        (
            "calibration",
            True,
            "How many vector dimensions are defined by VECTOR-2048?",
            ["05_architecture.pptx"],
        ),
        ("calibration", True, "What time does START-09 onboarding begin?", ["06_onboarding.pptx"]),
        (
            "calibration",
            True,
            "What is the acknowledgement time for P1-05?",
            ["07_incident_runbook.docx"],
        ),
        (
            "calibration",
            True,
            "How many remote days does REMOTE-3 allow?",
            ["08_remote_handbook.txt"],
        ),
        ("calibration", False, "What is the company helicopter registration number?", []),
        ("calibration", False, "Ai là giám đốc pháp lý của công ty?", []),
        (
            "calibration",
            True,
            "Nhân viên có chính xác bao nhiêu ngày phép năm?",
            ["01_leave_policy.pdf"],
        ),
        (
            "evaluation",
            True,
            "Theo thông báo scan, thẻ tạm thời hết hạn sau bao lâu?",
            ["02_security_scanned.pdf"],
        ),
        (
            "evaluation",
            True,
            "Mã WELL-12 cho biết bao nhiêu giờ wellness mỗi quý?",
            ["03_benefits_table.pdf"],
        ),
        ("evaluation", True, "Client meal code MEAL-40 has what limit?", ["04_finance_limits.csv"]),
        (
            "evaluation",
            True,
            "Which queue carries ingestion jobs asynchronously?",
            ["05_architecture.pptx"],
        ),
        ("evaluation", True, "Where do new staff collect equipment?", ["06_onboarding.pptx"]),
        ("evaluation", True, "LOG-90 stores audit logs for how long?", ["09_retention_policy.md"]),
        ("evaluation", True, "What is the reorder level for DOCK-12?", ["10_inventory.csv"]),
        ("evaluation", False, "What is the cafeteria Wi-Fi password?", []),
        ("evaluation", False, "Hãy cho biết giá cổ phiếu nội bộ ngày mai.", []),
    )
    query_variants = (
        "{query}",
        "Please answer only from the indexed corpus: {query}",
        "Return the exact documented value. {query}",
        "Internal policy lookup: {query}",
        "Retrieval verification question: {query}",
    )
    with path.open("x", encoding="utf-8") as handle:
        for index, (split, answerable, query, expected) in enumerate(base_samples, start=1):
            for variant_index, template in enumerate(query_variants, start=1):
                handle.write(
                    json.dumps(
                        {
                            "id": f"p4-{index:03d}-v{variant_index}",
                            "split": split,
                            "answerable": answerable,
                            "query": template.format(query=query),
                            "expected_source_names": expected,
                            "language": "vi" if any(ord(char) > 127 for char in query) else "en",
                            "tags": ["phase4-live-corpus", "retrieval"],
                        },
                        ensure_ascii=False,
                        sort_keys=True,
                    )
                    + "\n"
                )


async def _wait_job(
    client: httpx.AsyncClient,
    headers: Mapping[str, str],
    job_id: str,
    *,
    timeout_seconds: float,
) -> Mapping[str, object]:
    deadline = time.monotonic() + timeout_seconds
    last: Mapping[str, object] = {}
    while time.monotonic() < deadline:
        response = await client.get(f"/documents/jobs/{job_id}", headers=headers)
        if response.status_code != 200:
            raise RuntimeError(f"job status returned HTTP {response.status_code}")
        last = cast(Mapping[str, object], response.json())
        if last.get("state") in {"success", "failed", "cancelled"}:
            return last
        await asyncio.sleep(1)
    raise RuntimeError(f"job {job_id} timed out in state {last.get('state')}")


def _database_url(args: argparse.Namespace) -> URL:
    password = args.database_password_file.read_text(encoding="utf-8").strip()
    if not password:
        raise RuntimeError("database password file is empty")
    return URL.create(
        "postgresql+psycopg",
        username=args.database_user,
        password=password,
        host=args.database_host,
        port=args.database_port,
        database=args.database_name,
    )


def _broker_url(args: argparse.Namespace) -> str:
    password = args.broker_password_file.read_text(encoding="utf-8").strip()
    if not password:
        raise RuntimeError("broker password file is empty")
    return (
        f"amqp://{quote(args.broker_user, safe='')}:{quote(password, safe='')}@"
        f"{args.broker_host}:{args.broker_port}//"
    )


async def run(args: argparse.Namespace) -> int:
    output = args.output_dir.resolve()
    output.mkdir(parents=True, exist_ok=False)
    fixtures_dir = output / "fixtures"
    fixtures = generate_fixtures(fixtures_dir)
    if len(fixtures) != 10:
        raise RuntimeError("acceptance corpus must contain exactly 10 successful files")
    gold_path = output / "gold.jsonl"
    write_gold(gold_path)

    tenant_id = uuid5(NAMESPACE_URL, f"ntc-phase4-tenant:{output.name}")
    user_id = uuid5(NAMESPACE_URL, f"ntc-phase4-user:{output.name}")
    headers = {
        "X-NTC-Tenant-ID": str(tenant_id),
        "X-NTC-User-ID": str(user_id),
    }
    settings = IngestionSettings(
        enabled=True,
        trusted_actor_headers_enabled=True,
        database_host=args.database_host,
        database_port=args.database_port,
        database_name=args.database_name,
        database_user=args.database_user,
        database_password_file=args.database_password_file,
        minio_endpoint=args.minio_endpoint,
        minio_access_key=args.minio_access_key,
        minio_secret_key_file=args.minio_secret_key_file,
        minio_bucket=args.minio_bucket,
        broker_host=args.broker_host,
        broker_port=args.broker_port,
        broker_user=args.broker_user,
        broker_password_file=args.broker_password_file,
    )
    application = create_app(ingestion_settings=settings)
    uploads: list[dict[str, object]] = []
    transport = httpx.ASGITransport(app=application)
    async with application.router.lifespan_context(application):  # noqa: SIM117
        async with httpx.AsyncClient(transport=transport, base_url="http://phase4.local") as client:
            for fixture in fixtures:
                started = time.perf_counter()
                response = await client.post(
                    "/documents/upload",
                    headers=headers,
                    files={"file": (fixture.name, fixture.read_bytes())},
                )
                latency_ms = (time.perf_counter() - started) * 1000
                if response.status_code != 202:
                    raise RuntimeError(
                        f"upload {fixture.name} returned HTTP {response.status_code}: "
                        f"{response.text[:200]}"
                    )
                body = cast(dict[str, object], response.json())
                if body.get("duplicate") is not False:
                    raise RuntimeError(f"fresh upload {fixture.name} was unexpectedly duplicate")
                uploads.append({"source_name": fixture.name, "latency_ms": latency_ms, **body})

            jobs = await asyncio.gather(
                *(
                    _wait_job(
                        client,
                        headers,
                        cast(str, upload["job_id"]),
                        timeout_seconds=args.job_timeout_seconds,
                    )
                    for upload in uploads
                )
            )
            failures = [
                f"{upload['source_name']}: {job.get('error_code')} {job.get('error_message')}"
                for upload, job in zip(uploads, jobs, strict=True)
                if job.get("state") != "success"
            ]
            if failures:
                raise RuntimeError("successful corpus had worker failures: " + "; ".join(failures))

            duplicate_response = await client.post(
                "/documents/upload",
                headers=headers,
                files={"file": (fixtures[0].name, fixtures[0].read_bytes())},
            )
            duplicate_body = cast(dict[str, object], duplicate_response.json())
            if duplicate_response.status_code != 202 or duplicate_body.get("duplicate") is not True:
                raise RuntimeError("duplicate HTTP upload was not returned idempotently")

            document_id = cast(str, uploads[0]["document_id"])
            first_reindex = await client.post(f"/documents/{document_id}/reindex", headers=headers)
            second_reindex = await client.post(f"/documents/{document_id}/reindex", headers=headers)
            if first_reindex.status_code != 200 or second_reindex.status_code != 200:
                raise RuntimeError("reindex endpoint failed")
            first_reindex_body = cast(dict[str, object], first_reindex.json())
            second_reindex_body = cast(dict[str, object], second_reindex.json())
            if first_reindex_body["job_id"] != second_reindex_body["job_id"]:
                raise RuntimeError("concurrent reindex requests created duplicate jobs")
            reindex_job = await _wait_job(
                client,
                headers,
                cast(str, first_reindex_body["job_id"]),
                timeout_seconds=args.job_timeout_seconds,
            )
            if reindex_job.get("state") != "success":
                raise RuntimeError("reindex job did not succeed")

            corrupt = b"%PDF-1.7\n1 0 obj << /Type /Catalog >> broken-parser-fixture"
            corrupt_response = await client.post(
                "/documents/upload",
                headers=headers,
                files={"file": ("11_corrupt.pdf", corrupt)},
            )
            if corrupt_response.status_code != 202:
                raise RuntimeError("corrupt PDF did not reach asynchronous parser validation")
            corrupt_body = cast(dict[str, object], corrupt_response.json())
            corrupt_job = await _wait_job(
                client,
                headers,
                cast(str, corrupt_body["job_id"]),
                timeout_seconds=args.job_timeout_seconds,
            )
            if corrupt_job.get("state") != "failed" or not corrupt_job.get("error_code"):
                raise RuntimeError("corrupt PDF did not produce a useful failed job")
            for _ in range(2):
                deletion = await client.delete(
                    f"/documents/{corrupt_body['document_id']}", headers=headers
                )
                if deletion.status_code != 204:
                    raise RuntimeError("idempotent delete failed")

    engine = create_engine(_database_url(args), pool_pre_ping=True)
    try:
        with engine.connect() as connection:
            counts = {
                "documents": connection.execute(
                    text("""
                        SELECT COUNT(*) FROM documents
                        WHERE tenant_id=:tenant AND deleted_at IS NULL
                    """),
                    {"tenant": tenant_id},
                ).scalar_one(),
                "ready_documents": connection.execute(
                    text("""
                        SELECT COUNT(*) FROM documents
                        WHERE tenant_id=:tenant AND deleted_at IS NULL AND state='ready'
                    """),
                    {"tenant": tenant_id},
                ).scalar_one(),
                "versions": connection.execute(
                    text("""
                        SELECT COUNT(*) FROM document_versions v
                        JOIN documents d ON d.id=v.document_id
                        WHERE d.tenant_id=:tenant AND d.deleted_at IS NULL
                    """),
                    {"tenant": tenant_id},
                ).scalar_one(),
                "chunks": connection.execute(
                    text("""
                        SELECT COUNT(*) FROM chunks c
                        JOIN documents d ON d.id=c.document_id
                        WHERE d.tenant_id=:tenant AND d.deleted_at IS NULL
                    """),
                    {"tenant": tenant_id},
                ).scalar_one(),
                "current_child_chunks": connection.execute(
                    text("""
                        SELECT COUNT(*) FROM chunks c
                        JOIN documents d ON d.id=c.document_id
                        WHERE d.tenant_id=:tenant AND d.deleted_at IS NULL
                          AND c.version_id=d.current_version_id AND c.chunk_type='child'
                    """),
                    {"tenant": tenant_id},
                ).scalar_one(),
            }
            metadata = tuple(
                dict(row)
                for row in connection.execute(
                    text("""
                        SELECT d.source_name, c.page, c.slide, c.section_path,
                               c.text, c.content_hash, c.token_count
                        FROM chunks c JOIN documents d ON d.id=c.document_id
                        WHERE d.tenant_id=:tenant AND d.deleted_at IS NULL
                          AND c.version_id=d.current_version_id AND c.chunk_type='child'
                        ORDER BY d.source_name, c.chunk_index
                    """),
                    {"tenant": tenant_id},
                ).mappings()
            )
            duplicate_groups = connection.execute(
                text("""
                    SELECT COUNT(*) FROM (
                        SELECT version_id, chunk_type, chunk_index, COUNT(*)
                        FROM chunks c JOIN documents d ON d.id=c.document_id
                        WHERE d.tenant_id=:tenant
                        GROUP BY version_id, chunk_type, chunk_index HAVING COUNT(*) > 1
                    ) duplicates
                """),
                {"tenant": tenant_id},
            ).scalar_one()
            artifact_rows = tuple(
                dict(row)
                for row in connection.execute(
                    text("""
                        SELECT d.source_name, v.raw_artifact_path, v.normalized_artifact_path,
                               v.preview_artifact_path
                        FROM documents d JOIN document_versions v ON v.id=d.current_version_id
                        WHERE d.tenant_id=:tenant AND d.deleted_at IS NULL
                        ORDER BY d.source_name
                    """),
                    {"tenant": tenant_id},
                ).mappings()
            )
    finally:
        engine.dispose()

    if counts["documents"] != 10 or counts["ready_documents"] != 10:
        raise RuntimeError(f"expected 10 active ready documents, observed {counts}")
    if counts["versions"] != 11 or counts["current_child_chunks"] <= 0:
        raise RuntimeError(f"version/chunk acceptance failed: {counts}")
    if duplicate_groups != 0:
        raise RuntimeError("task rerun/reindex produced duplicate chunk coordinates")
    by_source: dict[str, list[Mapping[str, object]]] = {}
    for row in metadata:
        by_source.setdefault(cast(str, row["source_name"]), []).append(row)
    if not any(row["page"] is not None for row in by_source["01_leave_policy.pdf"]):
        raise RuntimeError("text PDF chunks lost page metadata")
    scanned_text = " ".join(cast(str, row["text"]) for row in by_source["02_security_scanned.pdf"])
    if "30" not in scanned_text:
        raise RuntimeError("scanned PDF OCR did not retain the acceptance fact")
    if not any(row["slide"] is not None for row in by_source["05_architecture.pptx"]):
        raise RuntimeError("PPTX chunks lost slide metadata")
    if not all(cast(str, row["content_hash"]) for row in metadata):
        raise RuntimeError("one or more chunks lost content hashes")

    minio_secret = args.minio_secret_key_file.read_text(encoding="utf-8").strip()
    minio = Minio(
        args.minio_endpoint,
        access_key=args.minio_access_key,
        secret_key=minio_secret,
        secure=False,
    )
    objects = tuple(
        item.object_name for item in minio.list_objects(args.minio_bucket, recursive=True)
    )
    artifact_values = {
        row[key]
        for row in artifact_rows
        for key in ("raw_artifact_path", "normalized_artifact_path", "preview_artifact_path")
    }
    if None in artifact_values:
        raise RuntimeError("raw/normalized/preview artifact path is missing")
    expected_paths = {cast(str, value).split("/", 1)[1] for value in artifact_values}
    if not expected_paths.issubset(set(objects)):
        raise RuntimeError("MinIO is missing one or more version-bound artifacts")

    first_job_id = UUID(cast(str, uploads[0]["job_id"]))
    before_replay_chunks = counts["chunks"]
    producer = create_celery(
        BrokerSettings(
            broker_url=SecretStr(_broker_url(args)),
            app_name="phase4-acceptance-replay",
        )
    )
    try:
        producer.send_task(
            "worker.tasks.process_document",
            args=(str(first_job_id),),
            task_id=str(uuid4()),
            queue="ingestion",
            routing_key="document.process",
        )
    finally:
        producer.close()
    await asyncio.sleep(args.replay_settle_seconds)
    replay_engine = create_engine(_database_url(args), pool_pre_ping=True)
    try:
        with replay_engine.connect() as connection:
            after_replay_chunks = connection.execute(
                text("""
                    SELECT COUNT(*) FROM chunks c JOIN documents d ON d.id=c.document_id
                    WHERE d.tenant_id=:tenant AND d.deleted_at IS NULL
                """),
                {"tenant": tenant_id},
            ).scalar_one()
    finally:
        replay_engine.dispose()
    if after_replay_chunks != before_replay_chunks:
        raise RuntimeError("redelivery of a completed job changed chunk count")

    report = {
        "schema_version": 1,
        "status": "PASS",
        "run_id": output.name,
        "completed_at": datetime.now(UTC).isoformat(),
        "tenant_id": str(tenant_id),
        "owner_id": str(user_id),
        "corpus": [
            {
                "source_name": fixture.name,
                "size_bytes": fixture.stat().st_size,
                "sha256": _sha256(fixture),
            }
            for fixture in fixtures
        ],
        "gold": {
            "path": "gold.jsonl",
            "sha256": _sha256(gold_path),
            "calibration_samples": 50,
            "evaluation_samples": 50,
        },
        "http_uploads": uploads,
        "jobs": list(jobs),
        "corrupt_parser_job": corrupt_job,
        "reindex_job": reindex_job,
        "database_counts": counts,
        "duplicate_chunk_coordinate_groups": duplicate_groups,
        "minio": {
            "bucket": args.minio_bucket,
            "object_count": len(objects),
            "all_current_version_artifacts_present": True,
        },
        "metadata_checks": {
            "text_pdf_page": True,
            "scanned_pdf_ocr": True,
            "pptx_slide": True,
            "section_path_stored": all("section_path" in row for row in metadata),
            "full_text_and_hash": True,
        },
        "idempotency": {
            "duplicate_upload_returned_existing": True,
            "concurrent_reindex_same_job": True,
            "completed_task_redelivery_chunk_count_before": before_replay_chunks,
            "completed_task_redelivery_chunk_count_after": after_replay_chunks,
            "delete_twice_http_204": True,
        },
        "acceptance": {
            "pdf_scanned_csv_pptx_e2e": True,
            "upload_non_blocking": max(cast(float, item["latency_ms"]) for item in uploads)
            < args.max_upload_latency_ms,
            "metadata_preserved": True,
            "rerun_no_duplicates": True,
            "useful_parser_error": True,
        },
    }
    if not all(report["acceptance"].values()):
        raise RuntimeError("one or more Phase 4 acceptance gates failed")
    report_json = json.dumps(report, ensure_ascii=False, indent=2, sort_keys=True) + "\n"
    _write_exclusive(output / "report.json", report_json)
    _write_exclusive(
        output / "report.md",
        "\n".join(
            (
                "# Phase 4 live ingestion acceptance",
                "",
                "- Status: **PASS**",
                f"- Run: `{output.name}`",
                f"- Successful documents: `{counts['ready_documents']}`",
                f"- Current child chunks: `{counts['current_child_chunks']}`",
                f"- MinIO objects observed: `{len(objects)}`",
                "- HTTP upload, Celery redelivery, reindex, parser failure, and delete "
                "idempotency: **PASS**",
                "",
            )
        ),
    )
    print(output / "report.json")
    return 0


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-dir", type=Path, required=True)
    parser.add_argument("--database-host", required=True)
    parser.add_argument("--database-port", type=int, default=5432)
    parser.add_argument("--database-name", default="ntc_rag")
    parser.add_argument("--database-user", default="ntc_app")
    parser.add_argument("--database-password-file", type=Path, required=True)
    parser.add_argument("--minio-endpoint", required=True)
    parser.add_argument("--minio-access-key", default="ntc_minio_admin")
    parser.add_argument("--minio-secret-key-file", type=Path, required=True)
    parser.add_argument("--minio-bucket", default="ntc-documents")
    parser.add_argument("--broker-host", required=True)
    parser.add_argument("--broker-port", type=int, default=5672)
    parser.add_argument("--broker-user", default="ntc_worker")
    parser.add_argument("--broker-password-file", type=Path, required=True)
    parser.add_argument("--job-timeout-seconds", type=float, default=900)
    parser.add_argument("--replay-settle-seconds", type=float, default=10)
    parser.add_argument("--max-upload-latency-ms", type=float, default=10_000)
    return parser.parse_args()


if __name__ == "__main__":
    raise SystemExit(asyncio.run(run(parse_args())))
